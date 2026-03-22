"""
Google Slides Output Module

Creates a Google Slides presentation from a PitchDeck model.
Uses the Google Slides API to create slides with proper layout,
typography, and speaker notes.

Requires a Google Cloud service account with Slides + Drive API enabled.
"""
from __future__ import annotations

import json
import os
from typing import Any

from loguru import logger

from core.models import PitchDeck, Slide, SlideType
from config import GOOGLE_CREDENTIALS_JSON, GOOGLE_SLIDES_SCOPES


# ── Google API Auth ───────────────────────────────────────────────────

def _get_slides_service():
    """Initialize and return Google Slides + Drive API services."""
    from google.oauth2 import service_account
    from googleapiclient.discovery import build

    # Load credentials from JSON string (for HF Spaces secrets)
    if GOOGLE_CREDENTIALS_JSON:
        creds_dict = json.loads(GOOGLE_CREDENTIALS_JSON)
        credentials = service_account.Credentials.from_service_account_info(
            creds_dict, scopes=GOOGLE_SLIDES_SCOPES
        )
    else:
        # Fallback: look for credentials file
        creds_file = os.getenv("GOOGLE_APPLICATION_CREDENTIALS", "credentials.json")
        credentials = service_account.Credentials.from_service_account_file(
            creds_file, scopes=GOOGLE_SLIDES_SCOPES
        )

    slides_service = build("slides", "v1", credentials=credentials)
    drive_service = build("drive", "v3", credentials=credentials)

    return slides_service, drive_service


# ── Layout Constants ──────────────────────────────────────────────────

# EMU (English Metric Units): 1 inch = 914400 EMU
EMU_INCH = 914_400
SLIDE_WIDTH = 10 * EMU_INCH    # 10 inches (widescreen)
SLIDE_HEIGHT = int(5.625 * EMU_INCH)  # 5.625 inches (16:9)

# Color palette (dark professional theme)
COLORS = {
    "bg_dark": {"red": 0.08, "green": 0.08, "blue": 0.12},
    "bg_accent": {"red": 0.12, "green": 0.12, "blue": 0.18},
    "text_white": {"red": 1.0, "green": 1.0, "blue": 1.0},
    "text_gray": {"red": 0.7, "green": 0.7, "blue": 0.75},
    "accent_blue": {"red": 0.2, "green": 0.5, "blue": 1.0},
    "accent_gold": {"red": 1.0, "green": 0.8, "blue": 0.2},
}

# Typography
FONT_TITLE = "Montserrat"
FONT_BODY = "Open Sans"


# ── Slide Creation Helpers ────────────────────────────────────────────

def _create_text_box(
    page_id: str,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    font_family: str = FONT_BODY,
    color: dict = None,
    bold: bool = False,
    alignment: str = "START",
) -> list[dict]:
    """Generate Slides API requests to create a text box."""
    if not text:
        return []

    color = color or COLORS["text_white"]
    element_id = f"{page_id}_text_{int(left)}_{int(top)}"

    requests = [
        # Create the shape
        {
            "createShape": {
                "objectId": element_id,
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": page_id,
                    "size": {
                        "width": {"magnitude": width * EMU_INCH, "unit": "EMU"},
                        "height": {"magnitude": height * EMU_INCH, "unit": "EMU"},
                    },
                    "transform": {
                        "scaleX": 1,
                        "scaleY": 1,
                        "translateX": left * EMU_INCH,
                        "translateY": top * EMU_INCH,
                        "unit": "EMU",
                    },
                },
            }
        },
        # Insert text
        {
            "insertText": {
                "objectId": element_id,
                "text": text,
                "insertionIndex": 0,
            }
        },
        # Style the text
        {
            "updateTextStyle": {
                "objectId": element_id,
                "style": {
                    "fontFamily": font_family,
                    "fontSize": {"magnitude": font_size, "unit": "PT"},
                    "foregroundColor": {"opaqueColor": {"rgbColor": color}},
                    "bold": bold,
                },
                "textRange": {"type": "ALL"},
                "fields": "fontFamily,fontSize,foregroundColor,bold",
            }
        },
        # Paragraph alignment
        {
            "updateParagraphStyle": {
                "objectId": element_id,
                "style": {"alignment": alignment},
                "textRange": {"type": "ALL"},
                "fields": "alignment",
            }
        },
    ]

    return requests


def _set_slide_background(page_id: str, color: dict = None) -> dict:
    """Set solid background color for a slide."""
    color = color or COLORS["bg_dark"]
    return {
        "updatePageProperties": {
            "objectId": page_id,
            "pageProperties": {
                "pageBackgroundFill": {
                    "solidFill": {
                        "color": {"rgbColor": color}
                    }
                }
            },
            "fields": "pageBackgroundFill.solidFill.color",
        }
    }


def _add_speaker_notes(page_id: str, notes_text: str) -> list[dict]:
    """Add speaker notes to a slide."""
    if not notes_text:
        return []
    return [
        {
            "insertText": {
                "objectId": f"{page_id}_notes",
                "text": notes_text,
                "insertionIndex": 0,
            }
        }
    ]


# ── Slide Layout Builders ────────────────────────────────────────────

def _build_title_slide(page_id: str, slide: Slide) -> list[dict]:
    """Title slide: Big title centered, subtitle below, genre line."""
    reqs = [_set_slide_background(page_id)]
    reqs += _create_text_box(
        page_id, slide.title,
        left=1, top=1.2, width=8, height=1.5,
        font_size=44, font_family=FONT_TITLE,
        color=COLORS["text_white"], bold=True, alignment="CENTER",
    )
    reqs += _create_text_box(
        page_id, slide.subtitle,
        left=1.5, top=2.8, width=7, height=1,
        font_size=20, color=COLORS["text_gray"], alignment="CENTER",
    )
    if slide.body:
        reqs += _create_text_box(
            page_id, slide.body,
            left=2, top=4, width=6, height=0.5,
            font_size=14, color=COLORS["accent_blue"], alignment="CENTER",
        )
    return reqs


def _build_content_slide(page_id: str, slide: Slide) -> list[dict]:
    """Standard content slide: title + subtitle + body/bullets."""
    reqs = [_set_slide_background(page_id)]

    # Title
    reqs += _create_text_box(
        page_id, slide.title,
        left=0.8, top=0.4, width=8.4, height=0.8,
        font_size=32, font_family=FONT_TITLE,
        color=COLORS["accent_gold"], bold=True,
    )

    # Subtitle
    if slide.subtitle:
        reqs += _create_text_box(
            page_id, slide.subtitle,
            left=0.8, top=1.3, width=8.4, height=0.8,
            font_size=18, color=COLORS["text_white"],
        )

    # Body text
    y_pos = 2.3 if slide.subtitle else 1.5
    if slide.body:
        reqs += _create_text_box(
            page_id, slide.body,
            left=0.8, top=y_pos, width=8.4, height=1.5,
            font_size=16, color=COLORS["text_gray"],
        )
        y_pos += 1.6

    # Bullets
    if slide.bullets:
        bullet_text = "\n".join(f"• {b}" for b in slide.bullets)
        reqs += _create_text_box(
            page_id, bullet_text,
            left=1.0, top=y_pos, width=8, height=2.5,
            font_size=15, color=COLORS["text_white"],
        )

    return reqs


def _build_closing_slide(page_id: str, slide: Slide) -> list[dict]:
    """Closing slide: centered title + CTA."""
    reqs = [_set_slide_background(page_id, COLORS["bg_accent"])]
    reqs += _create_text_box(
        page_id, slide.title,
        left=1, top=1.5, width=8, height=1,
        font_size=40, font_family=FONT_TITLE,
        color=COLORS["text_white"], bold=True, alignment="CENTER",
    )
    reqs += _create_text_box(
        page_id, slide.subtitle,
        left=2, top=2.8, width=6, height=0.6,
        font_size=24, color=COLORS["accent_gold"], alignment="CENTER",
    )
    if slide.body:
        reqs += _create_text_box(
            page_id, slide.body,
            left=1.5, top=3.8, width=7, height=1,
            font_size=16, color=COLORS["text_gray"], alignment="CENTER",
        )
    return reqs


LAYOUT_MAP = {
    SlideType.TITLE: _build_title_slide,
    SlideType.CLOSING: _build_closing_slide,
}


# ── Main Export Function ──────────────────────────────────────────────

def export_to_google_slides(
    deck: PitchDeck,
    share_with_email: str | None = None,
) -> dict[str, str]:
    """
    Create a Google Slides presentation from a PitchDeck.

    Args:
        deck: Complete PitchDeck with slides.
        share_with_email: Optional email to share the presentation with.

    Returns:
        Dict with 'presentation_id' and 'url'.
    """
    slides_service, drive_service = _get_slides_service()

    # 1. Create empty presentation
    presentation = slides_service.presentations().create(
        body={"title": deck.title}
    ).execute()

    pres_id = presentation["presentationId"]
    logger.info(f"Created presentation: {pres_id}")

    # 2. Get the default slide (Google always creates one)
    default_slides = presentation.get("slides", [])

    # 3. Build all slide requests
    all_requests = []

    for i, slide in enumerate(deck.slides):
        if i == 0 and default_slides:
            # Use the default first slide
            page_id = default_slides[0]["objectId"]
        else:
            # Create new blank slide
            page_id = f"slide_{i}"
            all_requests.append({
                "createSlide": {
                    "objectId": page_id,
                    "insertionIndex": i,
                    "slideLayoutReference": {"predefinedLayout": "BLANK"},
                }
            })

        # Apply layout based on slide type
        layout_fn = LAYOUT_MAP.get(slide.slide_type, _build_content_slide)
        all_requests.extend(layout_fn(page_id, slide))

    # 4. Execute all requests in batch
    if all_requests:
        slides_service.presentations().batchUpdate(
            presentationId=pres_id,
            body={"requests": all_requests},
        ).execute()
        logger.info(f"Applied {len(all_requests)} slide requests")

    # 5. Share if email provided
    url = f"https://docs.google.com/presentation/d/{pres_id}/edit"
    if share_with_email:
        try:
            drive_service.permissions().create(
                fileId=pres_id,
                body={
                    "type": "user",
                    "role": "writer",
                    "emailAddress": share_with_email,
                },
                sendNotificationEmail=False,
            ).execute()
            logger.info(f"Shared with {share_with_email}")
        except Exception as e:
            logger.warning(f"Failed to share: {e}")

    return {
        "presentation_id": pres_id,
        "url": url,
    }


# ── PowerPoint Fallback (no Google creds) ─────────────────────────────

def export_to_pptx(deck: PitchDeck, output_path: str = "pitch_deck.pptx") -> str:
    """
    Fallback: Export to PowerPoint using python-pptx.
    Works without any API keys.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in deck.slides:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(20, 20, 30)

        # Title
        left = Inches(1)
        top = Inches(0.5) if slide_data.slide_type != SlideType.TITLE else Inches(1.5)
        width = Inches(11)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.title
        p.font.size = Pt(36 if slide_data.slide_type == SlideType.TITLE else 28)
        p.font.color.rgb = RGBColor(255, 200, 50)
        p.font.bold = True
        if slide_data.slide_type in (SlideType.TITLE, SlideType.CLOSING):
            p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if slide_data.subtitle:
            top_sub = Inches(2.8) if slide_data.slide_type == SlideType.TITLE else Inches(1.6)
            txBox2 = slide.shapes.add_textbox(Inches(1), top_sub, Inches(11), Inches(0.8))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = slide_data.subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(255, 255, 255)
            if slide_data.slide_type in (SlideType.TITLE, SlideType.CLOSING):
                p2.alignment = PP_ALIGN.CENTER

        # Body
        if slide_data.body:
            txBox3 = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.5))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = slide_data.body
            p3.font.size = Pt(14)
            p3.font.color.rgb = RGBColor(180, 180, 190)

        # Bullets
        if slide_data.bullets:
            y = 4.2 if slide_data.body else 2.8
            txBox4 = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(10), Inches(2.5))
            tf4 = txBox4.text_frame
            tf4.word_wrap = True
            for j, bullet in enumerate(slide_data.bullets):
                if j == 0:
                    p4 = tf4.paragraphs[0]
                else:
                    p4 = tf4.add_paragraph()
                p4.text = f"• {bullet}"
                p4.font.size = Pt(14)
                p4.font.color.rgb = RGBColor(220, 220, 230)
                p4.space_after = Pt(8)

        # Speaker notes
        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.speaker_notes

    prs.save(output_path)
    logger.info(f"Saved PowerPoint to {output_path}")
    return output_path
